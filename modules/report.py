"""
PPTX 리포트 자동 생성 모듈
python-pptx + kaleido(차트 이미지 변환)를 사용하여
클라이언트 전달용 프레젠테이션을 자동 생성합니다.

차트 스타일은 대시보드(Toss-style)와 동일하게 적용됩니다.
"""
import io
from datetime import datetime

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import plotly.graph_objects as go

from modules.config import (
    BRAND_PRIMARY, CHART_COLORS, WMA_WEIGHTS,
)
from modules.data_processing import aggregate_metrics, calc_ctr, aggregate_by_weekday, media_month_stats


# ──────────────────────────────────────────────
# 색상 상수
# ──────────────────────────────────────────────
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK = RGBColor(0x19, 0x1F, 0x28)
_GRAY = RGBColor(0x8B, 0x95, 0xA1)
_ORANGE = RGBColor(0xF7, 0x93, 0x1D)
_BG_LIGHT = RGBColor(0xF4, 0xF5, 0xF7)

# PPTX용 차트 레이아웃 (Kaleido 렌더링 최적화)
_FONT = 'Arial, Helvetica, sans-serif'  # Kaleido에서 안정적 렌더링
_PPTX_LAYOUT = dict(
    template='plotly_white',
    font=dict(family=_FONT, size=13, color='#4E5968'),
    paper_bgcolor='white',
    plot_bgcolor='white',
    margin=dict(t=52, l=56, r=32, b=52),
    title=dict(text="", font=dict(size=16, color='#191F28', family=_FONT), x=0, xanchor='left'),
    xaxis=dict(showgrid=False, linecolor='#EBEEF2', tickfont=dict(size=12, color='#8B95A1'), title=''),
    yaxis=dict(gridcolor='#F0F1F3', gridwidth=1, showline=False, tickfont=dict(size=12, color='#8B95A1'), title=''),
    legend=dict(font=dict(size=12), bgcolor='rgba(0,0,0,0)', orientation='h', y=-0.18),
    showlegend=True,
)


# ──────────────────────────────────────────────
# 헬퍼
# ──────────────────────────────────────────────

def _set_slide_bg(slide, color: RGBColor = _WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text,
                 font_size=12, bold=False, color=_DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_rounded_rect(slide, left, top, width, height, fill_color=_WHITE, text="",
                       font_size=11, font_color=_DARK, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def _pptx_title(text: str) -> dict:
    """_PPTX_LAYOUT 타이틀 스타일 상속 + text 오버라이드"""
    d = dict(_PPTX_LAYOUT['title'])
    d['text'] = text
    return d


def _fig_to_image_bytes(fig: go.Figure, width=960, height=500) -> bytes:
    """Plotly Figure를 고해상도 PNG 바이트로 변환"""
    layout = {**_PPTX_LAYOUT}
    layout.update(width=width, height=height)
    fig.update_layout(**layout)
    return fig.to_image(format="png", scale=2)


def _add_chart_image(slide, fig: go.Figure, left=0.4, top=1.6, width=9.2, height=4.8):
    img_bytes = _fig_to_image_bytes(fig)
    slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        Inches(left), Inches(top), Inches(width), Inches(height),
    )


# ──────────────────────────────────────────────
# PPTX 전용 차트 빌더 (대시보드와 동일한 비주얼)
# ──────────────────────────────────────────────

def _chart_bar(df, x, y, title, color):
    """라운드 바 차트 — 대시보드 bar_chart() 동일 스타일"""
    _ymax = float(df[y].max()) * 1.18 if len(df) else 1.0
    fig = go.Figure(go.Bar(
        x=df[x], y=df[y],
        marker_color=color,
        marker_line_width=0,
        marker_cornerradius=6,
        opacity=0.88,
        text=df[y].apply(lambda v: f"{v:,.0f}"),
        textposition='outside',
        textfont=dict(size=11, color='#4E5968', family=_FONT),
        cliponaxis=False,
    ))
    fig.update_layout(
        bargap=0.35, showlegend=False,
        title=_pptx_title(title),
        yaxis=dict(range=[0, _ymax]),
        margin=dict(t=60, l=48, r=30, b=48),
    )
    return fig


def _chart_dual(df, x, bar_y, line_y, title):
    """이중축 바+라인 — 대시보드 dual_axis_bar_line() 동일 스타일"""
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df[x], y=df[bar_y], name="발송량",
        marker_color='#3182F6', marker_line_width=0, opacity=0.85,
        marker_cornerradius=6,
    ))
    fig.add_trace(go.Scatter(
        x=df[x], y=df[line_y], name="클릭수",
        mode='lines+markers', marker_color='#FF6B6B',
        line=dict(width=2.5, shape='spline'),
        marker=dict(size=8), yaxis='y2',
    ))
    fig.update_layout(
        title=_pptx_title(title),
        yaxis2=dict(overlaying='y', side='right', showgrid=False,
                    tickfont=dict(size=12, color='#8B95A1')),
    )
    return fig


def _chart_line(df, x, y, title, color):
    """스플라인 라인 차트 — 대시보드 line_chart() 동일 스타일"""
    fig = go.Figure(go.Scatter(
        x=df[x], y=df[y],
        mode='lines+markers',
        line=dict(color=color, width=2.5, shape='spline'),
        marker=dict(size=8, color=color, line=dict(width=1.5, color='white')),
        fill='tozeroy',
        fillcolor=f'rgba({int(color[1:3],16)},{int(color[3:5],16)},{int(color[5:7],16)},0.08)',
    ))
    fig.update_layout(showlegend=False,
                      title=_pptx_title(title))
    return fig


def _chart_media_prediction(df):
    """매체별 클릭 실적 + WMA 예측 — 대시보드 prediction 동일 스타일"""
    top5 = df.groupby('매체명')['발송량'].sum().nlargest(5).index.tolist()
    all_media = sorted(df['매체명'].unique())
    filtered = df[df['매체명'].isin(top5)]
    monthly = filtered.groupby(['년월', '매체명']).agg({'발송량': 'sum', '클릭수': 'sum'}).reset_index()
    monthly['CTR'] = calc_ctr(monthly)

    fig = go.Figure()
    w = np.array(WMA_WEIGHTS)

    for media in top5:
        m_df = monthly[monthly['매체명'] == media].sort_values('년월')
        if m_df.empty:
            continue
        color = CHART_COLORS[all_media.index(media) % len(CHART_COLORS)]

        # 실적 라인 (스플라인 + 마커)
        fig.add_trace(go.Scatter(
            x=m_df['년월'], y=m_df['클릭수'],
            mode='lines+markers', name=media,
            line=dict(color=color, width=2.5, shape='spline'),
            marker=dict(size=7, color=color, line=dict(width=1.5, color='white')),
        ))

        # WMA 예측 점선
        if len(m_df) >= 3:
            pred_ctr = float(np.dot(m_df['CTR'].tail(3).values, w))
            pred_vol = float(np.dot(m_df['발송량'].tail(3).values, w))
            pred_click = (pred_ctr / 100) * pred_vol
            last_ym = m_df['년월'].iloc[-1]
            last_dt = pd.to_datetime(f"{int(last_ym[:4])}-{int(last_ym[6:8]):02d}-01")
            next_lbl = (last_dt + pd.DateOffset(months=1)).strftime('%Y년 %m월')

            fig.add_trace(go.Scatter(
                x=[last_ym, next_lbl],
                y=[m_df['클릭수'].iloc[-1], pred_click],
                mode='lines+markers',
                line=dict(color=color, dash='dot', width=2),
                marker=dict(size=9, symbol='diamond', color=color,
                            line=dict(width=1.5, color='white')),
                showlegend=False,
            ))

    fig.update_layout(
        title=_pptx_title("매체별 클릭 예측 (WMA)"),
        hovermode="x unified",
    )
    return fig


# ──────────────────────────────────────────────
# 슬라이드 빌더
# ──────────────────────────────────────────────

def _slide_title(prs, company_name, df, service_name=""):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _DARK)

    date_range = f"{df['날짜'].min().strftime('%Y.%m.%d')} – {df['날짜'].max().strftime('%Y.%m.%d')}"
    generated = datetime.now().strftime('%Y년 %m월 %d일')

    _add_textbox(slide, 1.0, 2.0, 8.0, 0.8, company_name,
                 font_size=32, bold=True, color=_ORANGE)
    _add_textbox(slide, 1.0, 2.9, 8.0, 0.6, "LMS 성과 분석 리포트",
                 font_size=22, bold=True, color=_WHITE)
    _add_textbox(slide, 1.0, 3.6, 8.0, 0.4, f"분석 기간: {date_range}",
                 font_size=13, color=_GRAY)
    _add_textbox(slide, 1.0, 4.1, 8.0, 0.4, f"생성일: {generated}",
                 font_size=11, color=_GRAY)
    if service_name:
        _add_textbox(slide, 1.0, 6.2, 8.0, 0.3, f"Powered by {service_name}",
                     font_size=9, color=_GRAY, alignment=PP_ALIGN.LEFT)


def _slide_kpi(prs, df):
    """KPI 요약 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_LIGHT)
    _add_textbox(slide, 0.5, 0.4, 9.0, 0.5, "Executive Summary", font_size=22, bold=True)

    total_cost = df['집행금액'].sum()
    total_send = df['발송량'].sum()
    total_click = df['클릭수'].sum()
    total_ctr = (total_click / total_send * 100) if total_send > 0 else 0

    kpis = [
        ("총 집행금액", f"{total_cost:,.0f}원"),
        ("총 발송량", f"{total_send:,.0f}건"),
        ("총 클릭수", f"{total_click:,.0f}회"),
        ("평균 CTR", f"{total_ctr:.2f}%"),
    ]

    for i, (label, value) in enumerate(kpis):
        left = 0.4 + i * 2.35
        _add_rounded_rect(slide, left, 1.3, 2.15, 1.6, _WHITE)
        _add_textbox(slide, left + 0.15, 1.45, 1.85, 0.3, label, font_size=10, color=_GRAY)
        _add_textbox(slide, left + 0.15, 1.85, 1.85, 0.5, value,
                     font_size=20, bold=True, color=_DARK, alignment=PP_ALIGN.LEFT)

    # 월별 트렌드 (라운드 바)
    monthly = aggregate_metrics(df, ['년월'])
    fig = _chart_bar(monthly, '년월', '집행금액', "월별 집행금액 추이", BRAND_PRIMARY)
    _add_chart_image(slide, fig, left=0.4, top=3.2, width=9.2, height=3.8)


def _slide_trend(prs, df):
    """월별 발송량 & 클릭수 트렌드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    monthly = aggregate_metrics(df, ['년월'])
    fig = _chart_dual(monthly, '년월', '발송량', '클릭수', "월별 발송량 & 클릭수 트렌드")
    _add_chart_image(slide, fig, left=0.4, top=0.6, width=9.2, height=5.6)


def _slide_ctr_trend(prs, df):
    """CTR 트렌드 (스플라인 + 영역 채우기)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    monthly = aggregate_metrics(df, ['년월'])
    fig = _chart_line(monthly, '년월', 'CTR', "월별 CTR 추이", '#20C997')
    _add_chart_image(slide, fig, left=0.4, top=0.6, width=9.2, height=5.6)


def _slide_media_prediction(prs, df):
    """매체별 클릭 예측"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    fig = _chart_media_prediction(df)
    _add_chart_image(slide, fig, left=0.4, top=0.6, width=9.2, height=5.6)


def _slide_insights(prs, df):
    """인사이트 & 전략 제안"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_LIGHT)
    _add_textbox(slide, 0.5, 0.4, 9.0, 0.5, "전략 제안 및 Next Action", font_size=22, bold=True)

    insights = []

    # 요일별 효율
    day_stats = aggregate_by_weekday(df)
    best_day = day_stats.loc[day_stats['CTR'].idxmax()]
    worst_day = day_stats.loc[day_stats['CTR'].idxmin()]
    insights.append(
        f"📅  발송 타이밍: {best_day['짧은_요일']}요일 CTR {best_day['CTR']:.2f}%로 최고 / "
        f"{worst_day['짧은_요일']}요일 {worst_day['CTR']:.2f}%로 최저"
    )

    # 매체별 CTR 변동
    months = sorted(df['년월'].unique())
    if len(months) >= 2:
        latest = media_month_stats(df, months[-1]).set_index('매체명').rename(columns={'CTR': '최근CTR'})
        prev = media_month_stats(df, months[-2]).set_index('매체명').rename(columns={'CTR': '과거CTR'})
        trend = pd.merge(latest[['최근CTR']], prev[['과거CTR']], on='매체명').reset_index()
        trend['변화량'] = trend['최근CTR'] - trend['과거CTR']
        if not trend.empty:
            best = trend.loc[trend['변화량'].idxmax()]
            worst = trend.loc[trend['변화량'].idxmin()]
            insights.append(
                f"📈  효율 상승: {best['매체명']} CTR {best['변화량']:+.2f}%p → 예산 증액 추천"
            )
            insights.append(
                f"📉  효율 저하: {worst['매체명']} CTR {worst['변화량']:+.2f}%p → A/B 테스트 필요"
            )

    # 인사이트 모듈 연동
    try:
        from modules.insights import detect_summary
        cur = df[df['년월'] == months[-1]]
        prev_df = df[df['년월'] == months[-2]] if len(months) >= 2 else cur
        cur_s, prev_s = cur['발송량'].sum(), prev_df['발송량'].sum()
        cur_c, prev_c = cur['클릭수'].sum(), prev_df['클릭수'].sum()
        cur_ctr = (cur_c / cur_s * 100) if cur_s > 0 else 0
        prev_ctr = (prev_c / prev_s * 100) if prev_s > 0 else 0
        rule_insights = detect_summary(
            cur['집행금액'].sum(), prev_df['집행금액'].sum(),
            cur_s, prev_s, cur_c, prev_c, cur_ctr, prev_ctr, pd.DataFrame(),
        )
        for ri in rule_insights[:2]:
            insights.append(f"{ri['icon']}  {ri['fact']} → {ri['action']}")
    except Exception:
        pass

    insights.append(f"📊  총 {len(df['매체명'].unique())}개 매체, {len(months)}개월 분석 완료")

    for idx, text in enumerate(insights[:5]):
        _add_rounded_rect(slide, 0.5, 1.3 + idx * 1.05, 9.0, 0.85, _WHITE)
        _add_textbox(slide, 0.75, 1.42 + idx * 1.05, 8.5, 0.6, text,
                     font_size=13, color=_DARK, alignment=PP_ALIGN.LEFT)


# ──────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────

def generate_pptx(df: pd.DataFrame, company_name: str = "", service_name: str = "") -> bytes:
    """DataFrame을 기반으로 PPTX 리포트를 생성하고 바이트 반환"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_title(prs, company_name, df, service_name=service_name)
    _slide_kpi(prs, df)
    _slide_trend(prs, df)
    _slide_ctr_trend(prs, df)
    _slide_media_prediction(prs, df)
    _slide_insights(prs, df)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
